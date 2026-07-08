from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent
SRC = ROOT.parent / "2026-07-08-dia-dos-pais-stories-sqz" / "assets"
TMP = ROOT / "_prepared"
TMP.mkdir(parents=True, exist_ok=True)
OUT = ROOT / "sqz-stories-dia-dos-pais-editavel.pptx"

W, H = 7.5, 13.333333
WHITE = RGBColor(255, 255, 255)
OFF = RGBColor(232, 232, 232)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(8, 8, 8)

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

def crop_cover(src, name, ratio_w, ratio_h):
    im = Image.open(src).convert("RGB")
    sw, sh = im.size
    target = ratio_w / ratio_h
    current = sw / sh
    if current > target:
        nw = int(sh * target)
        left = (sw - nw) // 2
        box = (left, 0, left + nw, sh)
    else:
        nh = int(sw / target)
        top = (sh - nh) // 2
        box = (0, top, sw, top + nh)
    im = im.crop(box)
    out = TMP / name
    im.save(out, quality=92)
    return out

def add_bg(slide, color=DARK):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def add_overlay(slide, x, y, w, h, transparency=25):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = BLACK; r.fill.transparency = transparency
    r.line.fill.background()
    return r

def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=None, spacing=0.9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    p.line_spacing = spacing
    return box

def add_kicker(slide, text, page):
    add_text(slide, text.upper(), 0.48, 0.84, 4.8, 0.35, 12, True, OFF)
    add_text(slide, page, 6.32, 0.78, 0.7, 0.28, 9, True, OFF, PP_ALIGN.RIGHT)

def add_logo(slide):
    logo = SRC / "sqz-logo-horizontal-branco.png"
    slide.shapes.add_picture(str(logo), Inches(0.48), Inches(0.34), width=Inches(1.42))

def add_bar(slide, pct):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.48), Inches(12.72), Inches(6.54), Inches(0.025))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(85, 85, 85); bg.line.fill.background()
    fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.48), Inches(12.72), Inches(6.54*pct), Inches(0.025))
    fg.fill.solid(); fg.fill.fore_color.rgb = WHITE; fg.line.fill.background()

def add_tag(slide, text, x, y, w=1.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=text; p.alignment=PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name="Arial"; r.font.size=Pt(9); r.font.bold=True; r.font.color.rgb=BLACK

# Slide 1
s = prs.slides.add_slide(blank); add_bg(s)
s.shapes.add_picture(str(crop_cover(SRC/'bottle_jeep.jpg','s1_bottle.jpg',3.6,5.1)), Inches(3.7), Inches(1.45), width=Inches(3.55), height=Inches(5.1))
s.shapes.add_picture(str(crop_cover(SRC/'board_round.jpeg','s1_board.jpg',2.9,4.1)), Inches(0.48), Inches(5.35), width=Inches(2.9), height=Inches(4.1))
add_overlay(s,0,0,W,H,38); add_logo(s); add_kicker(s,'Dia dos Pais • SQZ','01/05')
add_text(s,'Presença que fica\nna rotina.',0.48,7.65,6.45,1.35,33,True)
add_text(s,'Produtos personalizados para reconhecer pais, equipes e parceiros com utilidade real.',0.51,9.05,6.1,0.72,13,False,OFF)
add_tag(s,'Story de abertura',4.95,10.85,1.95); add_bar(s,0.20)

# Slide 2
s = prs.slides.add_slide(blank); add_bg(s)
s.shapes.add_picture(str(crop_cover(SRC/'board_laser.jpeg','s2_board.jpg',6.54,5.9)), Inches(0.48), Inches(2.02), width=Inches(6.54), height=Inches(5.9))
add_overlay(s,0,0,W,H,32); add_logo(s); add_kicker(s,'Produto em foco','02/05')
add_text(s,'Kit churrasco\ntambém é estratégia.',0.48,1.65,6.45,1.35,30,True)
add_tag(s,'Gravação a laser',0.48,5.85,1.75)
card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.48), Inches(9.9), Inches(6.54), Inches(1.45))
card.fill.solid(); card.fill.fore_color.rgb=BLACK; card.fill.transparency=15; card.line.fill.background()
add_text(s,'Reconhecimento com uso real.',0.72,10.12,5.9,0.35,16,True)
add_text(s,'Tábua, utensílios e acessórios personalizados conectam marca, família e experiência fora do escritório.',0.72,10.57,5.8,0.6,11,False,OFF)
add_bar(s,0.40)

# Slide 3
s = prs.slides.add_slide(blank); add_bg(s, RGBColor(246,246,246))
s.shapes.add_picture(str(crop_cover(SRC/'bottle_jeep.jpg','s3_bottle.jpg',7.5,8.2)), Inches(0), Inches(1.8), width=Inches(7.5), height=Inches(8.2))
add_overlay(s,0,6.2,W,4.0,42); add_logo(s); add_kicker(s,'Uso diário','03/05')
add_text(s,'A marca acompanha\no dia.',0.48,8.98,6.55,1.12,30,True)
add_text(s,'Garrafas e squeezes funcionam porque entram na rotina: mesa, carro, treino, viagem e evento.',0.51,10.1,6.1,0.68,12,False,OFF)
add_tag(s,'Solução funcional',4.9,8.74,1.88); add_bar(s,0.60)

# Slide 4
s = prs.slides.add_slide(blank); add_bg(s)
s.shapes.add_picture(str(crop_cover(SRC/'backpack.jpeg','s4_bag.jpg',4.0,5.0)), Inches(3.0), Inches(1.5), width=Inches(4.0), height=Inches(5.0))
s.shapes.add_picture(str(crop_cover(SRC/'pen_desk.jpeg','s4_pen.jpg',2.7,3.4)), Inches(0.48), Inches(3.7), width=Inches(2.7), height=Inches(3.4))
add_overlay(s,0,0,W,H,36); add_logo(s); add_kicker(s,'Kit corporativo','04/05')
add_tag(s,'Curadoria SQZ',5.05,5.95,1.65)
add_text(s,'Do escritório ao\ndeslocamento.',0.48,7.9,6.55,1.15,30,True)
add_text(s,'Mochilas, ecobags, cadernos e canetas completam ações de relacionamento e endomarketing.',0.51,9.05,6.2,0.64,12,False,OFF)
add_bar(s,0.80)

# Slide 5
s = prs.slides.add_slide(blank); add_bg(s)
add_logo(s); add_kicker(s,'Próximo passo','05/05')
add_text(s,'Cada projeto é uma\nsolução — não um\npedido.',0.48,1.95,6.54,1.8,29,True)
add_text(s,'A SQZ monta a curadoria por objetivo, prazo, técnica e contexto de entrega.',0.51,3.78,6.1,0.55,12,False,OFF)
imgs = [
    ('board_round.jpeg',0.48,5.2,2.95,2.45), ('bottle_jeep.jpg',4.05,5.2,2.97,2.45),
    ('backpack.jpeg',0.48,7.92,2.15,2.95), ('pen_close.jpg',2.9,7.92,1.7,2.95), ('board_laser.jpeg',4.87,7.92,2.15,2.95),
]
for fname,x,y,w,h in imgs:
    p = crop_cover(SRC/fname, f's5_{fname}.jpg', w,h)
    s.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
add_overlay(s,0,0,W,H,20)
add_text(s,'Quer transformar Dia dos Pais em presença de marca?',0.48,11.15,6.4,0.55,15,True)
add_text(s,'Chama no direct e peça uma curadoria.',0.48,11.68,6.4,0.35,10,False,OFF)
add_bar(s,1.0)

prs.save(OUT)
print(OUT)
